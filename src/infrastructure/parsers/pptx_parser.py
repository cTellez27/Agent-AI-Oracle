import uuid
from pathlib import Path
from typing import Union, BinaryIO
from pptx import Presentation

from src.domain.entities.document import Document, DocumentMetadata
from src.domain.interfaces.parser_interface import BaseParserInterface


class PPTXParser(BaseParserInterface):
    """Extractor de texto para presentaciones de PowerPoint (.pptx)."""

    def supports_extension(self, extension: str) -> bool:
        return extension.lower() in [".pptx", ".ppt"]

    def parse(self, file_source: Union[str, Path, BinaryIO], file_name: str, category: str = "General") -> Document:
        try:
            prs = Presentation(file_source)
            slides_text = []

            for i, slide in enumerate(prs.slides):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_lines.append(shape.text.strip())

                if slide_lines:
                    slides_text.append(f"--- Diapositiva {i + 1} ---\n" + "\n".join(slide_lines))

            full_content = "\n\n".join(slides_text) if slides_text else "Presentación sin texto extraíble."

            size_bytes = 0
            if isinstance(file_source, (str, Path)):
                size_bytes = Path(file_source).stat().st_size
            elif hasattr(file_source, "seek") and hasattr(file_source, "tell"):
                file_source.seek(0, 2)
                size_bytes = file_source.tell()
                file_source.seek(0)

            metadata = DocumentMetadata(
                file_name=file_name,
                file_type="PPTX",
                file_size_bytes=size_bytes,
                category=category,
                additional_info={"num_slides": len(prs.slides)}
            )

            return Document(
                id=str(uuid.uuid4()),
                content=full_content,
                metadata=metadata
            )
        except Exception as e:
            raise ValueError(f"Error al procesar la presentación PowerPoint '{file_name}': {str(e)}")
